import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from io import StringIO
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import HuggingFaceInstructEmbeddings
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.prompts.prompt import PromptTemplate
from langchain.prompts import SystemMessagePromptTemplate
from langchain.prompts import HumanMessagePromptTemplate
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain, LLMChain
from htmlTemplates2 import new, bot_template, user_template, page_bg_img
from langchain.llms import HuggingFaceHub
from pylatexenc.latex2text import LatexNodes2Text
from langchain.schema import Document as docu
import os
import io
from pptx import Presentation
from docx import Document
import time
import concurrent.futures
from typing import List, Dict, Any
import threading
from functools import partial
from langchain.llms import HuggingFacePipeline
from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline

import logging
logging.basicConfig(level=logging.INFO)

class DocumentProcessor:
    @staticmethod
    def pdfread(pdf) -> List[str]:
        try:
            reader = PdfReader(pdf)
            return [page.extract_text() for page in reader.pages]
        except Exception as e:
            logging.error(f"Error processing PDF: {e}")
            st.error(f"Error processing PDF: {str(e)}")
            return []

    @staticmethod
    def readlat(lat: str) -> str:
        try:
            return LatexNodes2Text().latex_to_text(lat).strip()
        except Exception as e:
            logging.error(f"Error processing LaTeX: {e}")
            st.error(f"Error processing LaTeX: {str(e)}")
            return ""

    @staticmethod
    def readocx(poc: bytes) -> List[str]:
        try:
            with io.BytesIO(poc) as poi:
                document = Document(poi)
                return [p.text for p in document.paragraphs if p.text.strip()]
        except Exception as e:
            logging.error(f"Error processing DOCX: {e}")
            st.error(f"Error processing DOCX: {str(e)}")
            return []

    @staticmethod
    def readppt(pp: bytes) -> List[str]:
        try:
            with io.BytesIO(pp) as pio:
                prs = Presentation(pio)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text.extend(run.text for paragraph in shape.text_frame.paragraphs 
                                        for run in paragraph.runs if run.text.strip())
                return text
        except Exception as e:
            logging.error(f"Error processing PPT: {e}")
            st.error(f"Error processing PPT: {str(e)}")
            return []

class TypeWriter:
    def __init__(self, speed: int = 10):
        self.speed = speed
        self._stop_event = threading.Event()

    def stop(self):
        self._stop_event.set()

    def write(self, text: str, template: str):
        tokens = text.split()
        container = st.empty()
        for index in range(len(tokens) + 1):
            if self._stop_event.is_set():
                break
            curr_full_text = " ".join(tokens[:index])
            container.markdown(template.replace("{{MSG}}", curr_full_text), 
                               unsafe_allow_html=True)
            time.sleep(1 / self.speed)

class LuminAI:
    def __init__(self):
        self.typewriter = TypeWriter()
        self.doc_processor = DocumentProcessor()
        self.load_environment()
        self.initialize_session_state()
        self.setup_ui()

    def load_environment(self):
        load_dotenv()
        self.model_name = os.getenv("MODEL_NAME", "mistralai/Mistral-7B-Instruct-v0.1")
        self.embedding_model = os.getenv("EMBEDDING_MODEL", "hkunlp/instructor-xl")

    def initialize_session_state(self):
        initial_states = {
            "conversation": None,
            "chat_history": None,
            "downcontent": "",
            "my_text": "",  
            "processing": False,
        }
        for key, value in initial_states.items():
            if key not in st.session_state:
                st.session_state[key] = value

    def setup_ui(self):
        st.set_page_config(page_title="LuminAI Research Assistant", 
                           page_icon="🔬")
        st.markdown(page_bg_img, unsafe_allow_html=True)

    def get_text_chunks(self, text: str) -> List[str]:
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        return text_splitter.split_text(text)

    
       

    def get_vectorstore(self, text_chunks: List[str]):
        try:
            embeddings = HuggingFaceEmbeddings(model_name=self.embedding_model)
        
            return FAISS.from_texts(texts=text_chunks, embedding=embeddings)
        
    
        except Exception as e:
            logging.error(f"Error creating vector store: {e}")
            st.error(f"Error creating vector store: {str(e)}")
            return None   


    def get_conversation_chain(self, vectorstore):
        try:
            tokenizer = AutoTokenizer.from_pretrained("gpt2")
            model = AutoModelForCausalLM.from_pretrained("gpt2", device_map="auto")
            
            tokenizer.pad_token = tokenizer.eos_token
            
            
            hf_pipeline = pipeline(
                "text-generation",
                model=model,
                tokenizer=tokenizer,
                max_new_tokens=100,
                temperature=0.7,
                pad_token_id=tokenizer.eos_token_id
            )
            
            
            llm = HuggingFacePipeline(pipeline=hf_pipeline)
            
            
            memory = ConversationBufferMemory(
                memory_key="chat_history",
                return_messages=True,
                output_key="answer",  
                input_key="question"  
            )
            
            return ConversationalRetrievalChain.from_llm(
                llm=llm,
                retriever=vectorstore.as_retriever(search_kwargs={"k": 1}),
                memory=memory,
                return_source_documents=True,
                output_key="answer"  
            )
            
        except Exception as e:
            logging.error(f"Error: {str(e)}")
            st.error(f"Error: {str(e)}")
            return None
    
    

    def process_documents(self, docs) -> str:
        raw_text = ""
        lock = threading.Lock()
        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = []
            processors = {
                'pdf': self.doc_processor.pdfread,
                'docx': lambda doc: self.doc_processor.readocx(doc.read()),
                'pptx': lambda doc: self.doc_processor.readppt(doc.read()),
                'tex': lambda doc: self.doc_processor.readlat(doc.read().decode("utf-8")),
            }
            for doc in docs:
                ext = doc.name.split('.')[-1].lower()
                if ext in processors:
                    futures.append(executor.submit(processors[ext], doc))

            for future in concurrent.futures.as_completed(futures):
                try:
                    result = future.result()
                    with lock:
                        raw_text += ' '.join(result) if isinstance(result, list) else result
                except Exception as e:
                    logging.error(f"Error processing document: {e}")
                    st.error(f"Error processing document: {str(e)}")
        return raw_text

    def handle_userinput(self, user_question: str):
        try:
            response = st.session_state.conversation({'question': user_question})
            st.session_state.chat_history = response['chat_history']
            sources = response['source_documents']
            
            message = st.session_state.chat_history[-1]
            template = user_template if len(st.session_state.chat_history) % 2 == 1 else bot_template
            self.typewriter.write(message.content, template)

            if sources:
                st.subheader("Sources:")
                for source in sources:
                    st.markdown(source.page_content)

            st.session_state.downcontent += f"""
            Question: {user_question}
            Response: {message.content}
            
            Sources:
            {chr(10).join(s.page_content for s in sources)}
            
            {'=' * 50}
            """
        except Exception as e:
            logging.error(f"Error processing question: {e}")
            st.error(f"Error processing your question: {str(e)}")
    

    def main(self):
        st.title("LuminAI Research Assistant")
        st.header("Your Intelligent Research Companion")
        st.write(new, unsafe_allow_html=True)

        with st.sidebar:
            st.subheader("Settings")
            st.toggle("Enable typewriter effect", key="use_typewriter", value=True)
            st.toggle("Show processing time", key="show_processing_time", value=False)

        st.subheader("Ask me anything about your research documents")
        user_input = st.text_input("", key="widget", on_change=self.submit)
        
        if st.session_state.my_text:
            st.subheader("Your Question:")
            st.write(st.session_state.my_text)
            self.handle_userinput(st.session_state.my_text)

        st.subheader("Document Library")
        docs = st.file_uploader("Upload your research documents", 
                               accept_multiple_files=True,
                               type=['pdf', 'docx', 'pptx', 'tex'])

        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("Process Documents"):
                if not docs:
                    st.warning("Please upload at least one document.")
                    return
                
                start_time = time.time()
                with st.spinner("Processing your documents..."):
                    raw_text = self.process_documents(docs)
                    if raw_text:
                        text_chunks = self.get_text_chunks(raw_text)
                        vectorstore = self.get_vectorstore(text_chunks)
                        if vectorstore:
                            st.session_state.conversation = self.get_conversation_chain(vectorstore)
                            st.success("Documents processed successfully!")
                            if st.session_state.show_processing_time:
                                processing_time = time.time() - start_time
                                st.info(f"Processing time: {processing_time:.2f} seconds")
                        else:
                            st.error("Error creating vector store")
                    else:
                        st.error("No text could be extracted from the documents")

        with col3:
            if st.session_state.downcontent:
                st.download_button(
                    label="Download Chat History",
                    data=st.session_state.downcontent,
                    file_name="luminai_chat_history.txt"
                )

    @staticmethod
    def submit():
        st.session_state.my_text = st.session_state.widget
        st.session_state.widget = ""

if __name__ == '__main__':
    app = LuminAI()
    app.main()