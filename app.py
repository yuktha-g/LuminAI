import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates2 import new, bot_template, user_template, page_bg_img
from langchain_community.llms import CTransformers
from langchain.llms import HuggingFaceHub
import time
from pptx import Presentation
from docx import Document
import io
import os
from io import StringIO
from latex2text import LatexNodes2Text  # For LaTeX processing

# Function to read PDFs
def pdfread(pdf):
    reader = PdfReader(pdf)
    text = []
    for i in range(len(reader.pages)):
        text.append(reader.pages[i].extract_text())
    return text

# Function to read LaTeX files
def readlat(lat):
    te = LatexNodes2Text().latex_to_text(lat)
    te.strip()
    return te

# Function to read DOCX files
def readocx(poc):
    poi = io.BytesIO(poc)
    document = Document(poi)
    text = []
    for x in document.paragraphs:
        text.append(x.text)
    return text

# Function to read PPTX files
def readppt(pp):
    pio = io.BytesIO(pp)
    prs = Presentation(pio)
    text = []
    for slide in prs.slides:
        temp = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    temp.append(r.text)
        text.extend(temp)
    return text

# Typewriter effect for chat
def typewriter(text, template, speed):
    tokens = text.split()
    container = st.empty()
    for index in range(len(tokens) + 1):
        curr_full_text = " ".join(tokens[:index])
        container.markdown(template.replace("{{MSG}}", curr_full_text), unsafe_allow_html=True)
        time.sleep(1 / speed)

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

# Function to generate a vector store from text chunks
def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

# Function to initialize and get the conversation chain
def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

# Function to handle user input and display chat history
def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    
    for message in st.session_state.chat_history:
        if "🤖" in message.content:
            typewriter(message.content, bot_template, 10)
        else:
            typewriter(message.content, user_template, 10)

# Function to submit user text input
def submit():
    st.session_state.my_text = st.session_state.widget
    st.session_state.widget = ""

# Main function to set up the Streamlit app
def main():
    load_dotenv()
    st.set_page_config(page_title="LuminAI", page_icon=":star:")
    st.markdown(page_bg_img, unsafe_allow_html=True)

    # Initialize session state for conversation and chat history
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.title("LuminAI")
    st.write(new, unsafe_allow_html=True)

    # Text input for user question
    if "my_text" not in st.session_state:
        st.session_state.my_text = ""
    
    st.text_input("Ask your LuminAI assistant", key="widget", on_change=submit)
    user_question = st.session_state.my_text

    if user_question:
        handle_userinput(user_question)

    # Section to upload documents
    st.subheader("Document Library")
    docs = st.file_uploader("Upload your files (PDF, DOCX, PPTX, LaTeX) and click 'Process'", accept_multiple_files=True)
    
    if st.button("Process"):
        with st.spinner("Processing..."):
            raw_text = ""
            for doc in docs:
                ext = doc.name.split('.')[-1].lower()
                if ext == 'pdf':
                    raw_text += ' '.join(pdfread(doc))
                elif ext == 'docx':
                    raw_text += ' '.join(readocx(doc.read()))
                elif ext == 'pptx':
                    raw_text += ' '.join(readppt(doc.read()))
                elif ext == 'tex':
                    bd = doc.getvalue()
                    so = StringIO(bd.decode("utf-8"))
                    so = so.read()
                    raw_text += readlat(so)
            
            # Split text into chunks and generate vector store
            if raw_text:
                text_chunks = get_text_chunks(raw_text)
                vectorstore = get_vectorstore(text_chunks)
                st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()
